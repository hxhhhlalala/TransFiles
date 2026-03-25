#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
FP8 量化通信方案 · 技术分享 PPT（浅色主题，Wan2.2）
依赖: pip install python-pptx
运行: python gen_fp8_ppt.py
输出: fp8_quant_alltoall.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ──────────────────── 浅色主题色系 ────────────────────────
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BG       = RGBColor(0xF5, 0xF8, 0xFC)   # 极浅蓝灰背景
DARK     = RGBColor(0x0D, 0x20, 0x44)   # 深蓝（标题栏/重文字）
BLUE     = RGBColor(0x00, 0x6D, 0xCC)   # 主强调蓝
BLUE_L   = RGBColor(0xD6, 0xEA, 0xFF)   # 浅蓝色块
CYAN     = RGBColor(0x00, 0x99, 0xCC)   # 次强调
TEXT     = RGBColor(0x1A, 0x1A, 0x2E)   # 正文深色
SUBTEXT  = RGBColor(0x4A, 0x5A, 0x70)   # 次级文字灰
TINT1    = RGBColor(0xF0, 0xF5, 0xFF)   # 表格偶数行
TINT2    = WHITE                          # 表格奇数行
HDR_BG   = RGBColor(0x1E, 0x3A, 0x6E)   # 表头深蓝
HDR_FC   = RGBColor(0xAD, 0xD8, 0xFF)   # 表头青蓝字
GREEN    = RGBColor(0x16, 0x7B, 0x45)   # 正面结果绿
RED      = RGBColor(0xB9, 0x1C, 0x1C)   # 负面警示红
ORANGE   = RGBColor(0xC0, 0x5A, 0x00)   # 橙色提示
YELLOW   = RGBColor(0x92, 0x60, 0x00)   # 金黄提示（深背景下可读）
BORDER   = RGBColor(0xCC, 0xD9, 0xEA)   # 表格边框浅灰蓝

FONT     = "Microsoft YaHei"
SW       = Inches(13.33)
SH       = Inches(7.5)


# ──────────────────── 工具函数 ────────────────────────

def new_slide(prs, bg=BG):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return sl


def box(sl, x, y, w, h, fill=None, line=None, lw=Pt(0.75)):
    shp = sl.shapes.add_shape(1, x, y, w, h)
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line; shp.line.width = lw
    else:
        shp.line.fill.background()
    return shp


def txt(sl, text, x, y, w, h, size=Pt(13), bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = size; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color; r.font.name = FONT
    return tb


def txtml(sl, lines, x, y, w, h, size=Pt(12.5), color=TEXT,
          bold_idx=None, cmap=None, align=PP_ALIGN.LEFT, safter=Pt(4)):
    """多行文本，bold_idx=行索引集合，cmap={行idx:颜色}。"""
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = safter
        r = p.add_run(); r.text = line; r.font.size = size
        r.font.name = FONT; r.font.bold = (i in bold_idx) if bold_idx else False
        r.font.color.rgb = cmap.get(i, color) if cmap else color
    return tb


def hline(sl, x, y, w, color=BLUE, lw=Pt(1.5)):
    shp = sl.shapes.add_shape(1, x, y, w, Pt(2))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()


def _hex(rgb):
    return f'{int(rgb[0]):02X}{int(rgb[1]):02X}{int(rgb[2]):02X}'


def set_cell_style(cell, text, fill, fc=TEXT, size=Pt(12),
                   bold=False, align=PP_ALIGN.CENTER):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    sc = etree.SubElement(sf, qn('a:srgbClr')); sc.set('val', _hex(fill))
    cell.text = text
    p = cell.text_frame.paragraphs[0]; p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text; run.font.size = size
    run.font.name = FONT; run.font.color.rgb = fc; run.font.bold = bold


def make_table(sl, data, x, y, w, h, ratios,
               hdr_fill=HDR_BG, hdr_fc=HDR_FC,
               r1=TINT2, r2=TINT1, fc=TEXT,
               last_hl=False, last_fc=GREEN, last_fill=RGBColor(0xE6, 0xF4, 0xEC),
               font_size=Pt(12)):
    rows, cols = len(data), len(data[0])
    tbl = sl.shapes.add_table(rows, cols, x, y, w, h).table
    total = sum(ratios)
    for i, r in enumerate(ratios):
        tbl.columns[i].width = int(w * r / total)
    for ri, row in enumerate(data):
        if ri == 0:
            fill, f, bd = hdr_fill, hdr_fc, True
        elif last_hl and ri == rows - 1:
            fill, f, bd = last_fill, last_fc, True
        else:
            fill = r1 if ri % 2 == 1 else r2
            f, bd = fc, False
        for ci, val in enumerate(row):
            al = PP_ALIGN.LEFT if ci == 0 and ri > 0 else PP_ALIGN.CENTER
            set_cell_style(tbl.cell(ri, ci), str(val), fill, f,
                           size=font_size, bold=bd, align=al)
    return tbl


def slide_title(sl, title, subtitle=""):
    """统一的页面标题栏。"""
    box(sl, Inches(0), Inches(0), SW, Inches(0.88), fill=DARK)
    hline(sl, Inches(0), Inches(0.88), SW, color=BLUE, lw=Pt(3))
    txt(sl, title, Inches(0.4), Inches(0.1), Inches(9), Inches(0.68),
        size=Pt(26), bold=True, color=WHITE)
    if subtitle:
        txt(sl, subtitle, Inches(9.5), Inches(0.2), Inches(3.6), Inches(0.5),
            size=Pt(13), color=HDR_FC, align=PP_ALIGN.RIGHT)


def label_box(sl, label, x, y, w, h, bg=BLUE_L, fc=BLUE, size=Pt(14)):
    box(sl, x, y, w, h, fill=bg, line=BLUE, lw=Pt(0.75))
    txt(sl, label, x + Inches(0.12), y + Inches(0.05),
        w - Inches(0.15), h - Inches(0.1), size=size, bold=True, color=fc)


# ════════════════════════════════════════════════════════════════
#  SLIDE 1 · 封面
# ════════════════════════════════════════════════════════════════
def make_slide1(prs):
    sl = new_slide(prs, bg=WHITE)

    # 顶部深蓝色块
    box(sl, Inches(0), Inches(0), SW, Inches(2.2), fill=DARK)
    # 左侧蓝色竖条
    box(sl, Inches(0), Inches(2.2), Inches(0.12), SH - Inches(2.2), fill=BLUE)
    # 底部装饰条
    box(sl, Inches(0), SH - Inches(0.35), SW, Inches(0.35), fill=DARK)

    # 主标题
    txt(sl, "FP8 量化通信接入方案",
        Inches(0.5), Inches(0.3), Inches(10), Inches(1.2),
        size=Pt(44), bold=True, color=WHITE)

    # 蓝色分隔线 + 副标题
    hline(sl, Inches(0.5), Inches(2.45), Inches(9), color=BLUE, lw=Pt(2))
    txt(sl, "Ulysses 序列并行 · Ascend NPU · FP8 量化通信",
        Inches(0.5), Inches(2.62), Inches(10), Inches(0.65),
        size=Pt(22), color=CYAN)

    # 信息行
    txtml(sl,
        ["模型：Wan2.2（视频生成）",
         "序列并行度 SP = 4  ·  注意力头数 N = 40  ·  2026.03"],
        Inches(0.5), Inches(3.5), Inches(8.5), Inches(1.1),
        size=Pt(16), color=SUBTEXT, safter=Pt(6))

    # 右下角数据说明 badge
    box(sl, Inches(8.8), Inches(5.5), Inches(4.3), Inches(1.5),
        fill=TINT1, line=BORDER, lw=Pt(1))
    txtml(sl,
        ["Profiler 数据说明",
         "Before = BF16 BSND alltoall（基线）",
         "After  = FP8 BSND alltoall（代码确认：scatter_idx=2）"],
        Inches(8.95), Inches(5.6), Inches(4.1), Inches(1.3),
        size=Pt(11.5), color=SUBTEXT,
        cmap={0: BLUE}, bold_idx={0}, safter=Pt(1))


# ════════════════════════════════════════════════════════════════
#  SLIDE 2 · 方案背景 & 设计决策
# ════════════════════════════════════════════════════════════════
def make_slide2(prs):
    sl = new_slide(prs)
    slide_title(sl, "方案背景 & 设计决策", "Wan2.2 | SP=4 | N=40")

    # ── 左：动机 + BSND vs BNSD ──
    box(sl, Inches(0.25), Inches(1.05), Inches(6.1), Inches(6.3), fill=WHITE,
        line=BORDER, lw=Pt(1))

    label_box(sl, "▌ 为什么做 FP8 量化通信",
              Inches(0.38), Inches(1.18), Inches(5.8), Inches(0.42))
    txtml(sl,
        ["• Ulysses SP 每层 3 路 All-to-All（Q / K / V），通信量随 SP 线性增长",
         "• BF16 QKV A2A 实测 3,010 µs，是 attention 关键瓶颈",
         "• FP8 量化：2B/elem → 1B/elem，通信量减半",
         "• 配合 NPU FP8 FA kernel，attention 全段加速"],
        Inches(0.45), Inches(1.72), Inches(5.85), Inches(1.6),
        size=Pt(13), color=TEXT, safter=Pt(5))

    hline(sl, Inches(0.38), Inches(3.38), Inches(5.7), color=BORDER, lw=Pt(1))

    label_box(sl, "▌ 为什么选 BSND 方向",
              Inches(0.38), Inches(3.5), Inches(5.8), Inches(0.42))

    bsnd_data = [
        ["alltoall 方向", "内部 .contiguous() 次数", "结论"],
        ["BSND  scatter N → gather S", "2 次", "✓  选定方案"],
        ["BNSD  scatter S → gather N", "3 次（多 pre-A2A copy）", "✗  放弃"],
    ]
    make_table(sl, bsnd_data,
               Inches(0.38), Inches(4.02), Inches(5.88), Inches(1.1),
               ratios=[5, 3.5, 2],
               last_hl=True,
               last_fill=RGBColor(0xFF, 0xEB, 0xEB),
               last_fc=RED,
               font_size=Pt(12))

    txtml(sl,
        ["Q/K/V 各 1 路 A2A，BNSD 方向每路多 1 次大 tensor 内存拷贝",
         "→ 3 路 × 多 1 次 = 理论额外拷贝开销显著，代码已固定为 BSND 方向"],
        Inches(0.45), Inches(5.22), Inches(5.8), Inches(0.75),
        size=Pt(12), color=SUBTEXT, cmap={1: BLUE}, safter=Pt(3))

    # ── 右：旋转矩阵 + Scale 处理 ──
    box(sl, Inches(6.6), Inches(1.05), Inches(6.5), Inches(6.3), fill=WHITE,
        line=BORDER, lw=Pt(1))

    label_box(sl, "▌ 正交旋转矩阵（QuaRot）",
              Inches(6.73), Inches(1.18), Inches(6.2), Inches(0.42))
    txtml(sl,
        ["FP8 量化会截断尾部精度，但 Q/K 的 outlier 分布不均，",
         "直接量化导致注意力分数误差偏大。",
         "→ 量化前对 Q/K 左乘正交旋转矩阵 R：",
         "   • R 由 QR 分解生成（seed=42），各层各卡共享",
         "   • Q/K 共用同一 R，保证  Q·R·(K·R)ᵀ = QKᵀ  数值不变",
         "   • 旋转后数值分布更均匀，FP8 量化误差显著降低"],
        Inches(6.8), Inches(1.72), Inches(6.15), Inches(1.85),
        size=Pt(13), color=TEXT,
        cmap={2: BLUE, 3: BLUE, 4: BLUE, 5: GREEN}, safter=Pt(4))

    hline(sl, Inches(6.73), Inches(3.63), Inches(6.0), color=BORDER, lw=Pt(1))

    label_box(sl, "▌ Scale All-to-All & Trim",
              Inches(6.73), Inches(3.75), Inches(6.2), Inches(0.42))
    txtml(sl,
        ["fa_block_quant_preprocess 输出 scale 形状：",
         "  [B, N, ⌈S_local / bs⌉, 1]（始终 BNSD 风格）",
         "Scale alltoall 方向：scatter N(dim=1) → gather blocks(dim=2)",
         "alltoall 后：[B, N/P, P × ⌈S_local/bs⌉, 1]  →  需 trim！",
         "  原因：P × ⌈S_local/bs⌉ ≥ ⌈S_full/bs⌉，最多多 1 块",
         "  操作：scale[:, :, :⌈S_full/bs⌉, :].contiguous()",
         "不 trim 或不 .contiguous() → FA kernel 读到越界 scale → 数值错误"],
        Inches(6.8), Inches(4.28), Inches(6.15), Inches(2.7),
        size=Pt(12.5), color=TEXT,
        cmap={3: ORANGE, 6: RED},
        bold_idx={3, 6}, safter=Pt(4))


# ════════════════════════════════════════════════════════════════
#  SLIDE 3 · 非 Overlap 性能对比
# ════════════════════════════════════════════════════════════════
def make_slide3(prs):
    sl = new_slide(prs)
    slide_title(sl, "性能对比 · 非 Overlap 场景",
                "Before=BF16 BSND  |  After=FP8 BSND")

    # ── 左侧：详细对比表 ──
    box(sl, Inches(0.25), Inches(1.05), Inches(8.8), Inches(6.3), fill=WHITE,
        line=BORDER, lw=Pt(1))

    label_box(sl, "▌ 关键算子耗时对比（同一 Transformer 层，N=40，SP=4）",
              Inches(0.38), Inches(1.18), Inches(8.55), Inches(0.42))

    perf = [
        ["阶段", "Before（BF16 BSND）", "After（FP8 BSND）", "变化"],
        ["QKV A2A 耗时（3 次合计）",    "3,010 µs",  "1,507 µs",  "↓ 50%"],
        ["Scale A2A（3 次）",           "—",         "~15 µs",    "+15 µs"],
        ["per-head FP8 量化（×10 头）", "~1,310 µs", "0 µs",      "↓ 1,310 µs"],
        ["量化预处理（一次性，循环外）",  "—",         "~1,158 µs", "+1,158 µs"],
        ["post-A2A 转置（BSND→BNSD）",  "—",         "~243 µs",   "+243 µs"],
        ["通信+量化 净收益",             "",           "",          "≈ −1,397 µs"],
    ]
    make_table(sl, perf,
               Inches(0.38), Inches(1.72), Inches(8.6), Inches(3.6),
               ratios=[5, 3.5, 3.5, 2.5],
               last_hl=True,
               last_fill=RGBColor(0xE6, 0xF4, 0xEC),
               last_fc=GREEN,
               font_size=Pt(12.5))

    hline(sl, Inches(0.38), Inches(5.42), Inches(8.5), color=BORDER, lw=Pt(1))

    txtml(sl,
        ["① post-A2A 转置（+243µs）为 A2A 后 BSND→BNSD 格式转换，供 FA kernel 使用",
         "② 输出 All-to-All 两方案均为 BF16 BSND，耗时相同（~968 µs），不计入收益"],
        Inches(0.45), Inches(5.52), Inches(8.6), Inches(0.78),
        size=Pt(12), color=SUBTEXT, safter=Pt(4))

    # ── 右侧：A2A 条形 + 分析 ──
    box(sl, Inches(9.3), Inches(1.05), Inches(3.8), Inches(6.3), fill=WHITE,
        line=BORDER, lw=Pt(1))

    label_box(sl, "▌ A2A 耗时对比",
              Inches(9.45), Inches(1.18), Inches(3.5), Inches(0.42))

    bars = [
        ("BF16 Q A2A", 1044, ORANGE),
        ("FP8  Q A2A",  491, GREEN),
        ("BF16 K A2A",  978, ORANGE),
        ("FP8  K A2A",  524, GREEN),
        ("BF16 V A2A",  988, ORANGE),
        ("FP8  V A2A",  492, GREEN),
    ]
    max_v = 1100
    bw_total = Inches(1.8)
    bh = Inches(0.32)
    gap = Inches(0.15)
    bx = Inches(9.45)
    by = Inches(1.73)
    for lbl, val, clr in bars:
        bw = bw_total * val / max_v
        txt(sl, lbl, bx, by + Inches(0.04), Inches(1.2), bh, size=Pt(10.5), color=TEXT)
        box(sl, bx + Inches(1.22), by, bw, bh, fill=clr)
        txt(sl, f"{val}", bx + Inches(1.25) + bw, by + Inches(0.04),
            Inches(0.6), bh, size=Pt(10.5), color=SUBTEXT)
        by += bh + gap

    txtml(sl,
        ["FP8 通信量 ×½", "实测 A2A 降低约 50%", "与理论一致"],
        Inches(9.45), Inches(5.55), Inches(3.5), Inches(0.8),
        size=Pt(14), color=GREEN, bold_idx={0}, safter=Pt(3))


# ════════════════════════════════════════════════════════════════
#  SLIDE 4 · 核心数据流 + 接入 Checklist
# ════════════════════════════════════════════════════════════════
def make_slide4(prs):
    sl = new_slide(prs)
    slide_title(sl, "核心数据流  &  接入 Checklist", "FP8 BSND 最终方案")

    # ── 左：数据流 Before / After 并列 ──
    box(sl, Inches(0.25), Inches(1.05), Inches(6.5), Inches(6.3), fill=WHITE,
        line=BORDER, lw=Pt(1))
    label_box(sl, "▌ 数据流对比  Before vs After",
              Inches(0.38), Inches(1.18), Inches(6.2), Inches(0.36))

    # 子列标题
    CW = Inches(2.95)   # 单列宽
    LX = Inches(0.35)   # Before 起始 x
    RX = Inches(3.45)   # After  起始 x
    SZ = Pt(10)

    def col_hdr(sl, x, label, color):
        box(sl, x, Inches(1.65), CW, Inches(0.30), fill=color)
        txt(sl, label, x + Inches(0.05), Inches(1.67), CW - Inches(0.1), Inches(0.28),
            size=Pt(10.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    col_hdr(sl, LX, "Before  (BF16 A2A)", SUBTEXT)
    col_hdr(sl, RX, "After   (FP8 A2A)",  BLUE)

    # 分隔竖线
    box(sl, Inches(3.38), Inches(1.65), Inches(0.04), Inches(5.55), fill=BORDER)

    before_steps = [
        (WHITE,  DARK,   "输入 BSND [B,S_local,N,D] BF16"),
        (TINT1,  DARK,   "all_to_all_4D(BF16, s=2,g=1)\n→ BF16 [B,S_full,N/P,D]"),
        (TINT1,  DARK,   ".transpose(1,2).contiguous()\n→ BNSD  .split(1,dim=1)"),
        (BLUE_L, BLUE,   "for head_i in range(N/P):\n  Q/K × R\n  fa_block_quant_preprocess\n  npu_fused_infer_attn_score_v2\n  .transpose(1,2) → BSND"),
        (HDR_BG, WHITE,  ".cat(dim=2)\n→ BSND [B,S_full,N/P,D]"),
        (HDR_BG, WHITE,  "all_to_all_4D(BF16, s=1,g=2)\n→ BF16 [B,S_local,N,D]"),
    ]
    after_steps = [
        (WHITE,  DARK,   "输入 BSND [B,S_local,N,D] BF16"),
        (BLUE_L, BLUE,   "Q/K × R  (整体，循环外)\nfa_block_quant_preprocess(BSND)\n→ FP8 + scale"),
        (TINT1,  DARK,   "all_to_all_4D(FP8,   s=2,g=1) ↓50%\nall_to_all_4D(scale, s=1,g=2)+trim"),
        (TINT1,  DARK,   ".transpose(1,2).contiguous()\n→ BNSD  .split(1,dim=1)"),
        (BLUE_L, BLUE,   "for head_i in range(N/P):\n  npu_fused_infer_attn_score_v2\n  (quant_mode=7, dequant_scale_*)\n  .transpose(1,2) → BSND"),
        (HDR_BG, WHITE,  ".cat(dim=2)\n→ BSND [B,S_full,N/P,D]"),
        (HDR_BG, WHITE,  "all_to_all_4D(BF16, s=1,g=2)\n→ BF16 [B,S_local,N,D]"),
    ]

    def draw_col(sl, x, steps):
        sy = Inches(2.05)
        for fill, fc, label in steps:
            n_lines = label.count('\n') + 1
            sh = Inches(0.30) + Inches(0.18) * (n_lines - 1)
            box(sl, x, sy, CW, sh, fill=fill, line=BORDER, lw=Pt(0.4))
            txt(sl, label, x + Inches(0.06), sy + Inches(0.03),
                CW - Inches(0.1), sh - Inches(0.04),
                size=SZ, color=fc, wrap=True)
            sy += sh + Inches(0.04)
            if sy < Inches(7.15):
                txt(sl, "▼", x + CW / 2 - Inches(0.08), sy - Inches(0.04),
                    Inches(0.2), Inches(0.08), size=Pt(7), color=SUBTEXT,
                    align=PP_ALIGN.CENTER)

    draw_col(sl, LX, before_steps)
    draw_col(sl, RX, after_steps)

    # ── 右：Checklist ──
    box(sl, Inches(7.0), Inches(1.05), Inches(6.1), Inches(6.3), fill=WHITE,
        line=BORDER, lw=Pt(1))

    label_box(sl, "▌ 接入 Checklist",
              Inches(7.13), Inches(1.18), Inches(5.8), Inches(0.42))

    checks = [
        (GREEN,  True,  "正交旋转矩阵：seed=42，torch.linalg.qr"),
        (SUBTEXT,False, "    Q/K 共用同一矩阵，类级缓存，各卡一致"),
        (GREEN,  True,  "量化前保存 origin_dtype"),
        (SUBTEXT,False, "    FA out_dtype=origin_dtype（不能用 FP8 tensor 的 dtype）"),
        (GREEN,  True,  "fa_block_quant_preprocess(layout=\"BSND\")"),
        (SUBTEXT,False, "    Q block_size=128  /  K/V block_size=256"),
        (BLUE,   True,  "all_to_all_4D(FP8, scatter=2, gather=1)  ← BSND"),
        (BLUE,   True,  "Scale A2A(scatter=1, gather=2) + trim + .contiguous()"),
        (GREEN,  True,  "FA 前 .transpose(1,2) → BNSD；FA 后 .transpose(1,2) → BSND"),
        (GREEN,  True,  "FA kernel: npu_fused_infer_attention_score_v2"),
        (ORANGE, True,  "Overlap：FP8 A2A 需逐头异步触发（当前 BSND 批量方案有缺陷）"),
    ]
    cy = Inches(1.73)
    for clr, bold, text in checks:
        txt(sl, text, Inches(7.2), cy, Inches(5.75), Inches(0.45),
            size=Pt(12), bold=bold, color=clr)
        cy += Inches(0.46)


# ════════════════════════════════════════════════════════════════
#  SLIDE 5 · 附录：Overlap 场景
# ════════════════════════════════════════════════════════════════
def make_slide5(prs):
    sl = new_slide(prs)
    slide_title(sl, "【附录】Overlap 场景分析",
                "Before=BF16 per-head  |  After=FP8 BSND batch")

    # ── 左：机制对比表 ──
    box(sl, Inches(0.25), Inches(1.05), Inches(7.3), Inches(6.3), fill=WHITE,
        line=BORDER, lw=Pt(1))

    label_box(sl, "▌ 机制与性能对比",
              Inches(0.38), Inches(1.18), Inches(7.0), Inches(0.42))

    ov = [
        ["指标", "Before（BF16 per-head）", "After（FP8 BSND）"],
        ["A2A 总路数",       "30 路（逐头异步）",     "60 路（批量同步）"],
        ["A2A 触发时机",    "头 i A2A 与 FA[i-1] 并行", "33 路在 FA[0] 前同步完成\n27 路在 FA[0] 后同步完成"],
        ["FA[0] 前等待",    "首头 A2A 触发后 ~776µs", "全部 A2A 完成后 ~19,000µs"],
        ["FA[0] wait",     "211 µs",                "136 µs"],
        ["FA[1-9] wait 均值", "~168 µs  ✓",         "~530 µs  ✗"],
        ["FA 总耗时",        "~64,082 µs（BF16 FA）",  "~61,134 µs（FP8 FA 更快）"],
    ]
    make_table(sl, ov,
               Inches(0.38), Inches(1.72), Inches(7.1), Inches(3.3),
               ratios=[3, 4, 5], font_size=Pt(11.5))

    label_box(sl, "▌ 根因",
              Inches(0.38), Inches(5.1), Inches(7.0), Inches(0.38),
              bg=RGBColor(0xFF, 0xF3, 0xE0), fc=ORANGE)
    txtml(sl,
        ["FP8 A2A 批量同步执行：27 路（heads 1-9）在 FA[0] 后统一发射",
         "→ FA[1-9] 在全部 A2A 完成后才能开始，完全无流水线重叠",
         "FP8 FA kernel 本身更快（−2,948 µs），但串行化抵消收益",
         "改进方向：FP8 A2A 改为逐头异步触发（类 Before 的 AsStrided 模式）"],
        Inches(0.45), Inches(5.56), Inches(7.1), Inches(1.62),
        size=Pt(12.5), color=TEXT,
        cmap={1: RED, 3: BLUE}, bold_idx={1, 3}, safter=Pt(4))

    # ── 右：FA Wait 对比 ──
    box(sl, Inches(7.8), Inches(1.05), Inches(5.3), Inches(6.3), fill=WHITE,
        line=BORDER, lw=Pt(1))

    label_box(sl, "▌ 逐头 FA Wait 时间（µs）",
              Inches(7.93), Inches(1.18), Inches(5.0), Inches(0.42))

    wait_data = [
        ["头", "Before wait", "After wait", "Before dur.", "After dur."],
        ["FA[0]",  "211",  "136",  "4,719",  "4,843"],
        ["FA[1]",  "207",  "541",  "4,722",  "4,900"],
        ["FA[2]",  "144",  "622",  "4,764",  "4,921"],
        ["FA[3]",  "257",  "491",  "5,799",  "4,923"],
        ["FA[4]",  "223",  "498",  "5,718",  "5,382"],
        ["FA[5]",  "147",  "609",  "8,949",  "5,566"],
        ["FA[6]",  "155",  "513",  "7,721",  "5,985"],
        ["FA[7]",  "142",  "496",  "6,116",  "6,079"],
        ["FA[8]",  "134",  "477",  "5,988",  "8,134"],
        ["FA[9]",  "142",  "527",  "9,587",  "10,400"],
        ["合计/均值", "~176", "~491",  "64,082", "61,134"],
    ]
    make_table(sl, wait_data,
               Inches(7.93), Inches(1.72), Inches(5.1), Inches(4.6),
               ratios=[1.5, 2.5, 2.5, 2.5, 2.5],
               last_hl=True,
               last_fill=TINT1,
               last_fc=DARK,
               font_size=Pt(11))

    txtml(sl,
        ["wait：FA kernel 在队列等待上一 op 的调度延迟（非 A2A 等待）",
         "dur：FA kernel 实际执行时长"],
        Inches(7.93), Inches(6.42), Inches(5.1), Inches(0.7),
        size=Pt(10.5), color=SUBTEXT, safter=Pt(2))


# ════════════════════════════════════════════════════════════════
#  主程序
# ════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    make_slide1(prs)
    make_slide2(prs)
    make_slide3(prs)
    make_slide4(prs)
    make_slide5(prs)

    out = "fp8_quant_alltoall.pptx"
    prs.save(out)
    print(f"[OK] 已生成 {out}（5 页）")


if __name__ == "__main__":
    main()
